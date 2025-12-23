import tkinter as tk
from tkinter import ttk, messagebox, filedialog
import threading
import requests
from bs4 import BeautifulSoup, Comment
import re
import os
import time
import atexit
import math

# --- THƯ VIỆN SELENIUM (Điều khiển trình duyệt Chrome tự động) ---
from selenium import webdriver
from selenium.webdriver.chrome.service import Service
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from webdriver_manager.chrome import ChromeDriverManager

# --- THƯ VIỆN PYTHON-PPTX (Tạo file PowerPoint) ---
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR

class MassSlideApp:
    def __init__(self, root):
        self.root = root
        self.root.title("Hệ Thống Soạn Lễ (Regex Parser)")
        
        # --- CẤU HÌNH GIAO DIỆN CHÍNH ---
        # Đặt kích thước cửa sổ là 1000x650 theo yêu cầu
        self.root.geometry("1000x650")
        self.root.attributes('-topmost', 0)

        # --- KHỞI TẠO BIẾN DỮ LIỆU ---
        self.driver = None  # Biến giữ trình duyệt Chrome ẩn
        self.ktcgkpv_data = {}  # Biến lưu các bài đọc (Đáp ca, Tin mừng...) lấy từ web
        # Danh sách các phần bài hát cần soạn
        self.available_song_parts = ["Đức Mẹ", "Nhập Lễ", "Dâng Lễ", "Hiệp Lễ 1", "Hiệp Lễ 2", "Kết Lễ"]
        self.selected_queue = []  # Danh sách các phần người dùng ĐÃ CHỌN để soạn
        self.current_step_index = 0  # Chỉ số bước hiện tại (đang soạn phần nào)
        self.collected_songs = {}  # Lưu trữ nội dung bài hát đã lấy được
        self.search_results = []  # Lưu kết quả tìm kiếm Google
        self.current_parsed_song = {}  # Lưu bài hát đang xử lý hiện tại (đã tách khổ)
        self.checkbox_vars = {}  # Quản lý các checkbox chọn khổ hát
        self.mass_parts_vars = {} # Quản lý checkbox chọn phần lễ (Nhập lễ, Kết lễ...)
        
        self.is_first_search = True  # Cờ đánh dấu lần tìm kiếm đầu tiên (để hiện thông báo Captcha)
        self.is_searching = False    # Cờ trạng thái đang tìm kiếm

        # --- CHẠY CÁC TÁC VỤ NGẦM (BACKGROUND) ---
        self.init_driver_background()     # Khởi động Chrome ngầm
        self.fetch_ktcgkpv_background()   # Tải trước các bài đọc Lời Chúa
        self.setup_screen_selection()     # Hiển thị màn hình chọn phần lễ
        
        # Đăng ký hàm dọn dẹp (tắt Chrome) khi tắt phần mềm
        atexit.register(self.cleanup_driver)

    def cleanup_driver(self):
        """Hàm tắt trình duyệt Chrome khi thoát ứng dụng"""
        if self.driver:
            try: self.driver.quit()
            except: pass

    # --- CÁC WORKER CHẠY NGẦM (THREADING) ---
    def init_driver_background(self):
        """Tạo luồng riêng để bật Chrome, tránh làm đơ giao diện"""
        threading.Thread(target=self._init_driver_worker, daemon=True).start()

    def _init_driver_worker(self):
        """Cấu hình và khởi động Selenium Chrome Driver"""
        try:
            options = webdriver.ChromeOptions()
            options.add_argument("--disable-gpu") # Tắt tăng tốc phần cứng để ổn định hơn
            options.add_experimental_option('excludeSwitches', ['enable-logging']) # Tắt log rác
            options.page_load_strategy = 'eager' # Chế độ tải trang nhanh (không chờ ảnh load hết)
            
            # Tự động tải và cài đặt ChromeDriver phù hợp với phiên bản Chrome máy tính
            self.driver = webdriver.Chrome(service=Service(ChromeDriverManager().install()), options=options)
            self.driver.get("https://www.google.com") # Mở sẵn Google
            self.driver.minimize_window() # Thu nhỏ cửa sổ xuống taskbar
        except: pass

    def fetch_ktcgkpv_background(self):
        """Tạo luồng riêng để tải bài đọc Lời Chúa từ trang ktcgkpv.org"""
        threading.Thread(target=self._fetch_ktcgkpv_worker, daemon=True).start()

    def _fetch_ktcgkpv_worker(self):
        """Logic cào dữ liệu (Scraping) trang ktcgkpv.org"""
        url = "https://ktcgkpv.org/readings/mass-reading"
        # Các CSS Selector để định vị vị trí văn bản trên trang web
        selectors = {
            "Ca Nhập Lễ": "div.introit.reading.division div[class*='body']",
            "Đáp Ca": "div.responsory.division div[class*='body']", 
            "Tung Hô Tin Mừng": "div.gospel-acclam div[class*='body']",
            "Ca Hiệp Lễ": "div.communion.reading.division div[class*='body']"
        }
        try:
            # Dùng Chrome ẩn (headless) để tải cho nhanh, không hiện cửa sổ
            op = webdriver.ChromeOptions(); op.add_argument("--headless"); op.add_argument("--disable-gpu")
            tmp = webdriver.Chrome(service=Service(ChromeDriverManager().install()), options=op)
            tmp.get(url)
            
            # Chờ tối đa 10s để dữ liệu hiện ra
            try: WebDriverWait(tmp, 10).until(EC.presence_of_element_located((By.CSS_SELECTOR, "div.responsory.division")))
            except: pass
            time.sleep(2)
            
            # Duyệt qua từng phần và lấy text
            for k, sel in selectors.items():
                try: 
                    elem = tmp.find_element(By.CSS_SELECTOR, sel)
                    raw_text = elem.text
                    # Xử lý riêng cho Đáp Ca (bỏ chữ Đ. đầu dòng)
                    if k == "Đáp Ca":
                        first_line = raw_text.split("\n")[0]
                        clean_text = first_line.replace("Đ.", "").replace("Đ .", "").strip()
                        # Nếu dòng 1 quá ngắn (lỗi định dạng), lấy dòng 2
                        if len(clean_text) < 5 and len(raw_text.split("\n")) > 1:
                            clean_text = raw_text.split("\n")[1].strip()
                        self.ktcgkpv_data[k] = clean_text
                    else:
                        self.ktcgkpv_data[k] = raw_text
                except: self.ktcgkpv_data[k] = ""
            tmp.quit() # Tắt Chrome tạm sau khi xong
        except: pass

    # =========================================================================
    # CẤU HÌNH GIAO DIỆN (UI SETUP)
    # =========================================================================
    def setup_screen_selection(self):
        """Màn hình 1: Chọn các phần lễ muốn soạn"""
        for w in self.root.winfo_children(): w.destroy() # Xóa sạch màn hình cũ
        
        f = tk.Frame(self.root); f.place(relx=0.5, rely=0.5, anchor="center")
        tk.Label(f, text="CẤU TRÚC THÁNH LỄ", font=("Arial",20,"bold"), fg="#0D47A1").pack(pady=20)
        
        fc = tk.Frame(f); fc.pack(pady=10)
        self.mass_parts_vars = {}
        # Tạo danh sách checkbox
        for part in self.available_song_parts:
            var = tk.BooleanVar(value=True) 
            self.mass_parts_vars[part] = var
            tk.Checkbutton(fc, text=part, variable=var, font=("Arial",12)).pack(anchor="w", pady=2)
            
        tk.Button(f, text="BẮT ĐẦU SOẠN ➤", command=self.transition_to_editor, bg="#D32F2F", fg="white", font=("Arial",14,"bold"), padx=20, pady=10).pack(pady=20)

        # --- YÊU CẦU: Thêm dòng chữ bản quyền vào góc trái dưới ---
        tk.Label(self.root, text="Developed by Nguyen Tien Dung", font=("Arial", 9, "italic"), fg="gray").place(relx=0.0, rely=1.0, x=80, y=-8, anchor="sw")

    def transition_to_editor(self):
        """Chuyển đổi từ màn hình chọn sang màn hình soạn thảo"""
        self.selected_queue = []
        for p in self.available_song_parts:
            if self.mass_parts_vars[p].get(): self.selected_queue.append(p)
        self.collected_songs = {}
        
        # Nếu không chọn bài hát nào -> Chế độ chỉ làm bài đọc
        if not self.selected_queue:
            self.setup_editor_ui()
            self.lbl_status.config(text="CHẾ ĐỘ CHỈ TẢI BÀI ĐỌC (READING ONLY)", fg="blue")
            self.entry_kw.config(state="disabled"); self.btn_search.config(state="disabled")
            self.btn_fetch.config(state="disabled"); self.btn_next.config(text="💾 XUẤT FILE PPTX NGAY", command=self.generate_final_pptx, bg="#4CAF50")
            self.btn_skip.pack_forget()
            return

        self.current_step_index = 0
        self.setup_editor_ui()
        self.load_step_logic()

    def setup_editor_ui(self):
        """Màn hình 2: Giao diện chính (Tìm kiếm, Kết quả, Chọn đoạn)"""
        for w in self.root.winfo_children(): w.destroy()
        
        # Thanh trạng thái trên cùng
        self.lbl_status = tk.Label(self.root, text="...", font=("Arial",16,"bold"), fg="#D32F2F", bg="#E3F2FD", pady=10)
        self.lbl_status.pack(fill="x")
        
        body = tk.Frame(self.root, padx=10, pady=10); body.pack(fill="both", expand=True)
        
        # Khu vực 1: Tìm kiếm
        f1 = tk.LabelFrame(body, text="1. Tìm kiếm (Google)", font=("Arial",10,"bold")); f1.pack(fill="x", pady=5)
        self.entry_kw = tk.Entry(f1, font=("Arial",11))
        self.entry_kw.pack(side="left", fill="x", expand=True, padx=5, pady=5)
        self.entry_kw.bind("<Return>", lambda e: self.on_click_search()) # Enter để tìm
        self.btn_search = tk.Button(f1, text="Tìm kiếm", command=self.on_click_search, bg="#4CAF50", fg="white"); self.btn_search.pack(side="left", padx=5)

        # Khu vực 2: Kết quả tìm kiếm
        f2 = tk.LabelFrame(body, text="2. Kết quả", font=("Arial",10,"bold")); f2.pack(fill="x", pady=5)
        self.listbox = tk.Listbox(f2, height=5, font=("Arial",10))
        self.listbox.pack(side="left", fill="x", expand=True, padx=5, pady=5)
        self.btn_fetch = tk.Button(f2, text="Lấy nội dung", command=self.on_click_fetch, bg="#2196F3", fg="white"); self.btn_fetch.pack(side="left", fill="y", padx=5, pady=5)

        # Khu vực 3: Chọn đoạn bài hát (Preview)
        f3 = tk.LabelFrame(body, text="3. Chọn đoạn", font=("Arial",10,"bold")); f3.pack(fill="both", expand=True, pady=5)
        self.canvas = tk.Canvas(f3); self.scroll_frame = tk.Frame(self.canvas)
        sb = tk.Scrollbar(f3, command=self.canvas.yview); self.canvas.configure(yscrollcommand=sb.set)
        self.canvas.create_window((0,0), window=self.scroll_frame, anchor="nw")
        self.scroll_frame.bind("<Configure>", lambda e: self.canvas.configure(scrollregion=self.canvas.bbox("all")))
        self.canvas.pack(side="left", fill="both", expand=True); sb.pack(side="right", fill="y")

        # Khu vực chân trang (Nút Tiếp tục / Bỏ qua)
        f_foot = tk.Frame(self.root, pady=10, bg="#eee"); f_foot.pack(side="bottom", fill="x")
        self.btn_skip = tk.Button(f_foot, text="Bỏ qua", command=self.on_click_skip); self.btn_skip.pack(side="left", padx=20)
        self.btn_next = tk.Button(f_foot, text="TIẾP TỤC >>", command=self.on_click_next, bg="#FF9800", fg="white", font=("Arial",12,"bold")); self.btn_next.pack(side="right", padx=20)

        # --- YÊU CẦU: Thêm dòng chữ bản quyền vào góc trái dưới (màn hình 2) ---
        tk.Label(self.root, text="Developed by Nguyen Tien Dung", font=("Arial", 9, "italic"), fg="gray").place(relx=0.0, rely=1.0, x=80, y=-8, anchor="sw")

    # --- ĐIỀU KHIỂN LUỒNG (LOGIC FLOW) ---
    def load_step_logic(self):
        """Tải dữ liệu cho bước hiện tại (Ví dụ: Chuyển từ Nhập lễ -> Dâng lễ)"""
        if self.current_step_index < len(self.selected_queue):
            part = self.selected_queue[self.current_step_index]
            self.lbl_status.config(text=f"ĐANG SOẠN: {part.upper()} ({self.current_step_index+1}/{len(self.selected_queue)})")
            
            # Reset giao diện cho phần mới
            self.entry_kw.delete(0, tk.END); self.listbox.delete(0, tk.END)
            self.search_results = []; self.current_parsed_song = {}
            for c in self.scroll_frame.winfo_children(): c.destroy()
            
            # Bật lại các nút chức năng
            self.entry_kw.focus(); self.entry_kw.config(state="normal")
            self.btn_search.config(state="normal"); self.btn_fetch.config(state="normal")
            self.btn_next.config(text=f"XÁC NHẬN '{part}' & TIẾP >>", command=self.on_click_next, bg="#FF9800")
            self.btn_skip.pack(side="left", padx=20)
        else:
            # Nếu đã hết các phần -> Chuyển sang chế độ Xuất file
            self.lbl_status.config(text="HOÀN TẤT!", fg="green")
            self.btn_next.config(text="💾 XUẤT FILE PPTX", command=self.generate_final_pptx, bg="#4CAF50")
            self.btn_skip.pack_forget()
            self.entry_kw.config(state="disabled"); self.btn_search.config(state="disabled"); self.btn_fetch.config(state="disabled")

    def on_click_next(self):
        """Lưu dữ liệu phần hiện tại và chuyển sang phần kế tiếp"""
        part = self.selected_queue[self.current_step_index]
        data = []
        # Chỉ lưu các đoạn (Phiên khúc/Điệp khúc) mà người dùng ĐÃ TICK chọn
        for k, item in self.current_parsed_song.items():
            if self.checkbox_vars.get(k) and self.checkbox_vars[k].get():
                data.append({"header": k, "label": item['label'], "content": item['content']})
        
        if data: 
            self.collected_songs[part] = data
        else:
            # Nếu người dùng bấm Tiếp mà không chọn đoạn nào -> Hỏi có muốn bỏ qua không
            if not messagebox.askyesno("Xác nhận", f"Bỏ qua phần '{part}'?"): return
        
        self.current_step_index += 1
        self.load_step_logic()

    def on_click_skip(self):
        """Bỏ qua phần hiện tại, không lưu gì cả"""
        self.current_step_index += 1
        self.load_step_logic()

    # --- LOGIC TÌM KIẾM (SEARCH) ---
    def on_click_search(self):
        """Xử lý sự kiện bấm nút Tìm kiếm"""
        kw = self.entry_kw.get().strip()
        if not kw: return
        if not self.driver: messagebox.showinfo("Lỗi", "Chrome đang mở..."); return
        
        self.root.config(cursor="watch"); self.btn_search.config(state="disabled"); self.entry_kw.config(state="disabled")
        
        if self.is_first_search:
            messagebox.showinfo("Lưu ý Lần Đầu", "Sắp mở Google.\nHãy giải Captcha (nếu có) trên Chrome.\nSau đó chờ phần mềm tự lấy kết quả.\n\nBấm OK để bắt đầu.")
            
        self.is_searching = True
        # Chạy tìm kiếm trong luồng riêng để không đơ giao diện
        threading.Thread(target=self._search_polling_worker, args=(kw,), daemon=True).start()

    def _search_polling_worker(self, keyword):
        """Gửi lệnh tìm kiếm lên Google và chờ lấy link từ thanhcaivietnam.net"""
        try:
            # Cú pháp tìm kiếm giới hạn trong trang thanhcaivietnam.net
            self.driver.get(f"https://www.google.com/search?q=site:thanhcavietnam.net {keyword}")
            
            if self.is_first_search: 
                self.driver.maximize_window(); self.is_first_search = False
            else: 
                self.driver.minimize_window()
            
            max_retries = 120; found_links = []
            # Vòng lặp chờ kết quả (tối đa 60 giây)
            for _ in range(max_retries):
                if not self.is_searching: break
                try:
                    elements = self.driver.find_elements(By.CSS_SELECTOR, "a")
                    temp = []
                    for l in elements:
                        h = l.get_attribute("href"); t = l.text
                        # Lọc chỉ lấy link từ thanhcaivietnam.net
                        if h and "thanhcavietnam.net" in h and "google" not in h and t.strip():
                            temp.append({"title": t, "href": h})
                    if temp: found_links = temp; break 
                except: pass
                time.sleep(0.5)
            
            # Lọc trùng lặp
            self.search_results = []
            seen = set()
            for item in found_links:
                if item['href'] not in seen: self.search_results.append(item); seen.add(item['href'])
            self.root.after(0, self._finish_search)
        except: self.root.after(0, self._finish_search)

    def _finish_search(self):
        """Hiển thị kết quả tìm kiếm lên Listbox"""
        self.listbox.delete(0, tk.END)
        for item in self.search_results: self.listbox.insert(tk.END, item['title'])
        self.root.config(cursor=""); self.btn_search.config(state="normal"); self.entry_kw.config(state="normal"); self.entry_kw.focus()
        try: self.driver.minimize_window()
        except: pass
        self.root.deiconify(); self.root.lift(); self.root.focus_force()

    # --- LOGIC LẤY & LÀM SẠCH NỘI DUNG (FETCH & CLEAN) ---
    def on_click_fetch(self):
        """Bắt sự kiện bấm nút Lấy nội dung"""
        sel = self.listbox.curselection()
        if not sel: return
        url = self.search_results[sel[0]]['href']
        self.root.config(cursor="watch")
        threading.Thread(target=self._fetch_worker, args=(url,), daemon=True).start()

    def _fetch_worker(self, url):
        """
        Hàm quan trọng nhất: Tải HTML, xóa rác (PDF, MP3, View more...)
        và tách văn bản thô để chuẩn bị xử lý.
        """
        try:
            headers = {'User-Agent': 'Mozilla/5.0'}
            resp = requests.get(url, headers=headers)
            soup = BeautifulSoup(resp.text, 'html.parser')

            # --- BƯỚC 1: DỌN RÁC CƠ BẢN (SCRIPT, STYLE, QUẢNG CÁO) ---
            for comment in soup.find_all(string=lambda text: isinstance(text, Comment)):
                comment.extract()
            for tag in soup.find_all(["fieldset", "script", "style", "iframe", "object", "h2", "h3", "h4"]):
                tag.decompose()

            # Xóa các thẻ div bị ẩn (display:none)
            for div in soup.find_all("div", style=True):
                if "display:none" in div.get("style", "").replace(" ", "").lower():
                    div.decompose()

            # --- BƯỚC 2: XÁC ĐỊNH VÙNG NỘI DUNG CHÍNH ---
            # Web này thường để nội dung trong thẻ div có id bắt đầu bằng 'post_message_'
            content_div = soup.find("div", id=re.compile(r"^post_message_"))
            
            if content_div:
                # Xóa TABLE đầu bài (nếu có)
                for table in content_div.find_all("table"):
                    table.decompose()

                # Xóa Metadata (Tác giả, Lời, Nhạc...)
                meta_keywords = ["Sáng tác", "Tác giả", "Lời:", "Nhạc:", "Ý:", "Thơ:", "Imprimatur"]
                for tag in content_div.find_all(["div", "p", "span", "b", "strong", "i", "font"]):
                    if any(k in tag.get_text() for k in meta_keywords):
                        tag.decompose()

                # Xóa các thành phần căn giữa (thường là Tiêu đề hoặc Nút Download PDF/MP3)
                center_pattern = re.compile(r"text-align\s*:\s*center", re.IGNORECASE)
                center_tags = content_div.find_all(lambda t:
                    (t.has_attr("align") and "center" in t["align"].lower()) or
                    (t.has_attr("style") and center_pattern.search(t["style"]))
                )
                for tag in center_tags:
                    # Kiểm tra kỹ: nếu thẻ center này chứa PDF/MP3/ENCORE thì xóa nó VÀ CẮT ĐUÔI luôn
                    txt_upper = tag.get_text().upper()
                    if "PDF" in txt_upper or "MP3" in txt_upper or "ENCORE" in txt_upper:
                        # Đây là điểm cắt! Xóa thẻ này và toàn bộ các thẻ đi sau nó (siblings)
                        current = tag
                        while current:
                            next_sibling = current.find_next_sibling()
                            while next_sibling:
                                next_sibling.decompose() # Xóa sạch các thẻ em phía sau (bao gồm cả View more...)
                                next_sibling = current.find_next_sibling()
                            current = None # Thoát vòng lặp
                        tag.decompose() # Xóa chính thẻ nút bấm
                    else:
                        tag.decompose() # Xóa thẻ căn giữa bình thường (tiêu đề bài hát)

                # --- BƯỚC 3: XỬ LÝ CỤ THỂ "VIEW MORE" (PHÒNG HỜ) ---
                # Tìm thẻ nào chứa chữ "View more the latest threads" và xóa nó
                trash_markers = ["View more the latest threads", "Các chủ đề tương tự", "Copyright"]
                for marker in trash_markers:
                    found_node = content_div.find(string=re.compile(marker, re.IGNORECASE))
                    if found_node:
                        parent = found_node.parent
                        if parent and parent.name != 'div': parent = parent.parent
                        if parent:
                            for sib in list(parent.find_next_siblings()): sib.decompose()
                            parent.decompose()

                # --- BƯỚC 4: LẤY TEXT SẠCH ---
                q = content_div.find("blockquote")
                target = q if q else content_div
                
                # Lấy text. Mặc định separator='\n'
                txt = target.get_text(separator='\n')

                # --- BƯỚC 5: CẮT CHUỖI LẦN CUỐI (SAFETY NET) ---
                # Phòng trường hợp xóa HTML sót, ta cắt bằng cách duyệt từng dòng
                lines = txt.split('\n')
                clean_lines = []
                stop_keywords = ["VIEW MORE", "CÁC CHỦ ĐỀ TƯƠNG TỰ", "PDF", "MP3", "ENCORE", "LINK FILE", "CHIPLOVE"]
                
                for line in lines:
                    l_upper = line.upper().strip()
                    # Nếu gặp dòng chứa từ khóa dừng -> Dừng hẳn luôn
                    if any(k == l_upper for k in stop_keywords) or "VIEW MORE THE LATEST THREADS" in l_upper:
                        break
                    # Nếu dòng quá ngắn mà chứa từ khóa file -> Dừng
                    if len(l_upper) < 20 and any(k in l_upper for k in ["PDF", "MP3", "ENCORE"]):
                        break
                    clean_lines.append(line)

                final_text = "\n".join(clean_lines)

                # Gửi text đã làm sạch sang hàm phân tích
                self._parse_text(final_text)
                self.root.after(0, self._update_parsed_ui)
            else:
                self.root.after(0, lambda: messagebox.showerror("Lỗi", "Không tìm thấy nội dung bài hát (post_message div)."))

        except Exception as e:
            self.root.after(0, lambda: messagebox.showerror("Lỗi", str(e)))
        finally:
            self.root.after(0, lambda: self.root.config(cursor=""))

    def _parse_text(self, text):
        """
        Phân tích văn bản thô thành các object (Điệp khúc, Phiên khúc).
        Sử dụng Regex để nhận diện.
        """
        lines = [l.strip() for l in text.split('\n') if l.strip()]
        self.current_parsed_song = {}
        
        # Regex Patterns (Các mẫu nhận diện)
        p_chorus = re.compile(r'^(ĐK|DK|Đk|Dk|Điệp Khúc|Diep Khuc|Chorus)[:\.\s]?', re.IGNORECASE)
        p_verse = re.compile(r'^(\d+)([\.\)\/\-\:])') # Nhận diện số đầu dòng: 1. , 2. , 1) ...
        p_end = re.compile(r'^(CODA|KẾT|KEÁT|FINE)', re.IGNORECASE)

        current_header = "Mở đầu"
        current_label = ""
        current_content = []

        ignore_lines = ["PDF", "MP3", "ENCORE", "PDF MP3"]

        def save_section(h, l, c):
            """Hàm phụ để lưu đoạn vừa đọc được vào từ điển"""
            if c:
                key = h
                cnt = 1
                # Nếu trùng tên (ví dụ 2 đoạn Điệp Khúc) thì đánh số thêm
                while key in self.current_parsed_song:
                    key = f"{h} ({cnt})"
                    cnt += 1
                self.current_parsed_song[key] = {"label": l, "content": "\n".join(c)}

        for line in lines:
            if line.upper() in ignore_lines: continue
            
            match_chorus = p_chorus.match(line)
            match_verse = p_verse.match(line)
            match_end = p_end.match(line)
            
            is_new_section = False
            
            # Nếu phát hiện dòng bắt đầu bằng ĐK -> Là phần mới
            if match_chorus:
                is_new_section = True
                new_header = "Điệp Khúc"
                new_label = "ĐK:"
                content_start_idx = match_chorus.end()
            
            # Nếu phát hiện dòng bắt đầu bằng Số -> Là Phiên khúc mới
            elif match_verse:
                is_new_section = True
                number = match_verse.group(1) 
                new_header = f"Phiên Khúc {number}"
                new_label = f"{number}."
                content_start_idx = match_verse.end()
            
            # Nếu phát hiện Kết/Coda
            elif match_end:
                is_new_section = True
                new_header = line.strip()
                new_label = line.strip()
                content_start_idx = len(line)
            
            # Nếu là dòng in hoa ngắn (tiêu đề phụ)
            elif len(line) < 40 and line.isupper() and len(line) > 3 and not any(c in line for c in ",."):
                 is_new_section = True
                 new_header = line.title()
                 new_label = line
                 content_start_idx = len(line)

            if is_new_section:
                # Lưu đoạn trước đó lại
                save_section(current_header, current_label, current_content)
                # Bắt đầu đoạn mới
                current_header = new_header
                current_label = new_label
                current_content = []
                # Lấy phần chữ còn lại sau dấu hiệu nhận biết
                remaining_text = line[content_start_idx:].strip()
                while remaining_text and remaining_text[0] in [".", ":", " ", ")", "/"]:
                    remaining_text = remaining_text[1:].strip()
                if remaining_text:
                    current_content.append(remaining_text)
            else:
                current_content.append(line)
        
        # Lưu đoạn cuối cùng
        save_section(current_header, current_label, current_content)

    def _update_parsed_ui(self):
        """Cập nhật giao diện: Vẽ lại danh sách các checkbox chọn đoạn"""
        for c in self.scroll_frame.winfo_children(): c.destroy()
        self.checkbox_vars = {}
        for k, v in self.current_parsed_song.items():
            var = tk.BooleanVar(value=True); self.checkbox_vars[k] = var
            r = tk.Frame(self.scroll_frame); r.pack(fill="x", pady=2)
            tk.Checkbutton(r, text=k, variable=var, font=("Arial",10,"bold"), fg="#0D47A1").pack(side="left")
            tk.Label(r, text=f"[{v['label']}] {v['content'][:40]}...", fg="gray").pack(side="left")

    # --- TẠO FILE PPTX (GENERATE) ---
    def generate_final_pptx(self):
        """Hàm tạo file PowerPoint cuối cùng từ dữ liệu đã thu thập"""
        path = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint", "*.pptx")])
        if not path: return
        
        prs = Presentation()
        # Thiết lập kích thước slide 16:9
        prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

        def add_black():
            """Thêm một slide màu đen (ngăn cách các phần)"""
            s = prs.slides.add_slide(prs.slide_layouts[6])
            s.background.fill.solid(); s.background.fill.fore_color.rgb = RGBColor(0,0,0)

        def has_song(key): return key in self.collected_songs
        def has_read(key): return bool(self.ktcgkpv_data.get(key))

        def add_song(key):
            """Hàm thêm các slide bài hát (Xử lý logic lặp lại Điệp Khúc)"""
            if not has_song(key): return
            add_black() # Ngăn cách bằng slide đen
            secs = self.collected_songs[key]
            
            # 1. Tìm trước dữ liệu Điệp Khúc để dùng chèn sau
            chorus_data = next((s for s in secs if "Điệp Khúc" in s['header'] or "ĐK" in s['header']), None)

            for s in secs:
                # 2. Nếu đoạn hiện tại là Điệp Khúc -> BỎ QUA (để tránh in lặp)
                if "Điệp Khúc" in s['header'] or "ĐK" in s['header']:
                    continue

                # 3. Tạo slide cho Phiên Khúc
                self.create_slide(prs, key.upper(), s['label'], s['content'])

                # 4. Ngay sau Phiên Khúc -> Chèn luôn Điệp Khúc (nếu có)
                if "Phiên Khúc" in s['header'] and chorus_data:
                    self.create_slide(prs, key.upper(), chorus_data['label'], chorus_data['content'])

        def add_read(name):
            """Hàm thêm slide bài đọc (Đáp ca, Tin mừng...)"""
            txt = self.ktcgkpv_data.get(name, "")
            if not txt: return
            add_black()
            self.create_slide(prs, name.upper(), "", txt)

        # --- THỰC HIỆN TẠO SLIDE THEO THỨ TỰ PHỤNG VỤ ---
        if has_song("Đức Mẹ"): add_song("Đức Mẹ")
        if has_song("Nhập Lễ"): add_song("Nhập Lễ")
        if has_read("Ca Nhập Lễ"): add_read("Ca Nhập Lễ")

        if has_read("Đáp Ca") or has_read("Tung Hô Tin Mừng"):
            add_black()
            if has_read("Đáp Ca"): self.create_slide(prs, "ĐÁP CA", "", self.ktcgkpv_data["Đáp Ca"])
            if has_read("Tung Hô Tin Mừng"): self.create_slide(prs, "TUNG HÔ TIN MỪNG", "", self.ktcgkpv_data["Tung Hô Tin Mừng"])

        if has_song("Dâng Lễ"): add_song("Dâng Lễ")
        if has_read("Ca Hiệp Lễ"): add_read("Ca Hiệp Lễ")
        if has_song("Hiệp Lễ 1"): add_song("Hiệp Lễ 1")
        if has_song("Hiệp Lễ 2"): add_song("Hiệp Lễ 2")
        if has_song("Kết Lễ"): add_song("Kết Lễ")

        try:
            prs.save(path)
            messagebox.showinfo("Xong", f"Đã lưu: {path}")
            os.startfile(path) # Tự động mở file sau khi lưu
        except Exception as e: messagebox.showerror("Lỗi", str(e))

    def create_slide(self, prs, title, label, content):
        """Hàm vẽ chi tiết 1 slide (Tiêu đề, Kẻ ngang, Nội dung)"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid(); slide.background.fill.fore_color.rgb = RGBColor(0, 32, 96) # Màu nền xanh đậm
        sw, sh = prs.slide_width, prs.slide_height

        # Tiêu đề (Trên cùng)
        tb = slide.shapes.add_textbox(Inches(0), Inches(0), sw, Inches(1.4))
        p = tb.text_frame.paragraphs[0]
        p.text = title; p.alignment = PP_ALIGN.CENTER
        p.font.name = "Times New Roman"; p.font.size = Pt(72); p.font.bold = True; p.font.color.rgb = RGBColor(255, 255, 0)

        # Kẻ ngang màu vàng
        ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0), Inches(1.5), sw, Inches(1.5))
        ln.line.color.rgb = RGBColor(255, 255, 0); ln.line.width = Pt(3)

        # Nội dung chính
        tb2 = slide.shapes.add_textbox(Inches(0), Inches(1.6), sw, sh - Inches(2.0))
        tf = tb2.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
        
        # Xử lý xuống dòng và cỡ chữ tự động
        raw = " ".join([l.strip() for l in content.split('\n') if l.strip()])
        length = len(raw)
        lines_54 = math.ceil(length / 32)
        lines_44 = math.ceil(length / 42) 
        if lines_54 <= 6: size = 54
        elif lines_44 <= 9: size = 44
        elif length < 700: size = 36
        else: size = 28

        p2 = tf.paragraphs[0]; p2.alignment = PP_ALIGN.JUSTIFY; p2.line_spacing = 1.1
        if label: # Nếu có nhãn (ĐK, 1, 2...) thì tô đỏ
            r = p2.add_run(); r.text = label + " "
            r.font.name = "Times New Roman"; r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = RGBColor(255, 0, 0)
        
        # Phần lời hát màu trắng
        r2 = p2.add_run(); r2.text = raw
        r2.font.name = "Times New Roman"; r2.font.size = Pt(size); r2.font.bold = True; r2.font.color.rgb = RGBColor(255, 255, 255)

if __name__ == "__main__":
    root = tk.Tk()
    app = MassSlideApp(root)
    root.mainloop()