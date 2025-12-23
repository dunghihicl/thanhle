import tkinter as tk
from tkinter import ttk, messagebox, filedialog
import threading
import requests
from bs4 import BeautifulSoup, Comment
import re
import os
import time
import atexit
import math

# --- THƯ VIỆN SELENIUM ---
from selenium import webdriver
from selenium.webdriver.chrome.service import Service
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from webdriver_manager.chrome import ChromeDriverManager

# --- THƯ VIỆN PPTX ---
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR

class MassSlideApp:
    def __init__(self, root):
        self.root = root
        self.root.title("Hệ Thống Soạn Lễ (Smart & Clean)")
        self.root.geometry("1100x850")
        self.root.attributes('-topmost', 0)

        # Dữ liệu
        self.driver = None 
        self.ktcgkpv_data = {}
        self.available_song_parts = ["Đức Mẹ", "Nhập Lễ", "Dâng Lễ", "Hiệp Lễ 1", "Hiệp Lễ 2", "Kết Lễ"]
        self.selected_queue = [] 
        self.current_step_index = 0
        self.collected_songs = {} 
        self.search_results = []
        self.current_parsed_song = {} 
        self.checkbox_vars = {} 
        self.mass_parts_vars = {}
        
        self.is_first_search = True 
        self.is_searching = False

        self.init_driver_background()
        self.fetch_ktcgkpv_background()
        self.setup_screen_selection()
        atexit.register(self.cleanup_driver)

    def cleanup_driver(self):
        if self.driver:
            try: self.driver.quit()
            except: pass

    # --- BACKGROUND WORKERS ---
    def init_driver_background(self):
        threading.Thread(target=self._init_driver_worker, daemon=True).start()

    def _init_driver_worker(self):
        try:
            options = webdriver.ChromeOptions()
            options.add_argument("--disable-gpu")
            options.add_experimental_option('excludeSwitches', ['enable-logging'])
            options.page_load_strategy = 'eager'
            self.driver = webdriver.Chrome(service=Service(ChromeDriverManager().install()), options=options)
            self.driver.get("https://www.google.com")
            self.driver.minimize_window()
        except: pass

    def fetch_ktcgkpv_background(self):
        threading.Thread(target=self._fetch_ktcgkpv_worker, daemon=True).start()

    def _fetch_ktcgkpv_worker(self):
        url = "https://ktcgkpv.org/readings/mass-reading"
        selectors = {
            "Ca Nhập Lễ": "div.introit.reading.division div[class*='body']",
            "Đáp Ca": "div.responsory.division div[class*='body']", 
            "Tung Hô Tin Mừng": "div.gospel-acclam div[class*='body']",
            "Ca Hiệp Lễ": "div.communion.reading.division div[class*='body']"
        }
        try:
            op = webdriver.ChromeOptions(); op.add_argument("--headless"); op.add_argument("--disable-gpu")
            tmp = webdriver.Chrome(service=Service(ChromeDriverManager().install()), options=op)
            tmp.get(url)
            try: WebDriverWait(tmp, 10).until(EC.presence_of_element_located((By.CSS_SELECTOR, "div.responsory.division")))
            except: pass
            time.sleep(2)
            
            for k, sel in selectors.items():
                try: 
                    elem = tmp.find_element(By.CSS_SELECTOR, sel)
                    raw_text = elem.text
                    if k == "Đáp Ca":
                        first_line = raw_text.split("\n")[0]
                        clean_text = first_line.replace("Đ.", "").replace("Đ .", "").strip()
                        if len(clean_text) < 5 and len(raw_text.split("\n")) > 1:
                            clean_text = raw_text.split("\n")[1].strip()
                        self.ktcgkpv_data[k] = clean_text
                    else:
                        self.ktcgkpv_data[k] = raw_text
                except: self.ktcgkpv_data[k] = ""
            tmp.quit()
        except: pass

    # =========================================================================
    # UI SETUP
    # =========================================================================
    def setup_screen_selection(self):
        for w in self.root.winfo_children(): w.destroy()
        f = tk.Frame(self.root); f.place(relx=0.5, rely=0.5, anchor="center")
        tk.Label(f, text="CẤU TRÚC THÁNH LỄ", font=("Arial",20,"bold"), fg="#0D47A1").pack(pady=20)
        fc = tk.Frame(f); fc.pack(pady=10)
        self.mass_parts_vars = {}
        for part in self.available_song_parts:
            var = tk.BooleanVar(value=True) 
            self.mass_parts_vars[part] = var
            tk.Checkbutton(fc, text=part, variable=var, font=("Arial",12)).pack(anchor="w", pady=2)
        tk.Button(f, text="BẮT ĐẦU SOẠN ➤", command=self.transition_to_editor, bg="#D32F2F", fg="white", font=("Arial",14,"bold"), padx=20, pady=10).pack(pady=20)

    def transition_to_editor(self):
        self.selected_queue = []
        for p in self.available_song_parts:
            if self.mass_parts_vars[p].get(): self.selected_queue.append(p)
        self.collected_songs = {}
        
        if not self.selected_queue:
            self.setup_editor_ui()
            self.lbl_status.config(text="CHẾ ĐỘ CHỈ TẢI BÀI ĐỌC (READING ONLY)", fg="blue")
            self.entry_kw.config(state="disabled"); self.btn_search.config(state="disabled")
            self.btn_fetch.config(state="disabled"); self.btn_next.config(text="💾 XUẤT FILE PPTX NGAY", command=self.generate_final_pptx, bg="#4CAF50")
            self.btn_skip.pack_forget()
            return

        self.current_step_index = 0
        self.setup_editor_ui()
        self.load_step_logic()

    def setup_editor_ui(self):
        for w in self.root.winfo_children(): w.destroy()
        self.lbl_status = tk.Label(self.root, text="...", font=("Arial",16,"bold"), fg="#D32F2F", bg="#E3F2FD", pady=10)
        self.lbl_status.pack(fill="x")
        
        body = tk.Frame(self.root, padx=10, pady=10); body.pack(fill="both", expand=True)
        
        f1 = tk.LabelFrame(body, text="1. Tìm kiếm (Google)", font=("Arial",10,"bold")); f1.pack(fill="x", pady=5)
        self.entry_kw = tk.Entry(f1, font=("Arial",11))
        self.entry_kw.pack(side="left", fill="x", expand=True, padx=5, pady=5)
        self.entry_kw.bind("<Return>", lambda e: self.on_click_search())
        self.btn_search = tk.Button(f1, text="Tìm kiếm", command=self.on_click_search, bg="#4CAF50", fg="white"); self.btn_search.pack(side="left", padx=5)

        f2 = tk.LabelFrame(body, text="2. Kết quả", font=("Arial",10,"bold")); f2.pack(fill="x", pady=5)
        self.listbox = tk.Listbox(f2, height=5, font=("Arial",10))
        self.listbox.pack(side="left", fill="x", expand=True, padx=5, pady=5)
        self.btn_fetch = tk.Button(f2, text="Lấy nội dung", command=self.on_click_fetch, bg="#2196F3", fg="white"); self.btn_fetch.pack(side="left", fill="y", padx=5, pady=5)

        f3 = tk.LabelFrame(body, text="3. Chọn đoạn", font=("Arial",10,"bold")); f3.pack(fill="both", expand=True, pady=5)
        self.canvas = tk.Canvas(f3); self.scroll_frame = tk.Frame(self.canvas)
        sb = tk.Scrollbar(f3, command=self.canvas.yview); self.canvas.configure(yscrollcommand=sb.set)
        self.canvas.create_window((0,0), window=self.scroll_frame, anchor="nw")
        self.scroll_frame.bind("<Configure>", lambda e: self.canvas.configure(scrollregion=self.canvas.bbox("all")))
        self.canvas.pack(side="left", fill="both", expand=True); sb.pack(side="right", fill="y")

        f_foot = tk.Frame(self.root, pady=10, bg="#eee"); f_foot.pack(side="bottom", fill="x")
        self.btn_skip = tk.Button(f_foot, text="Bỏ qua", command=self.on_click_skip); self.btn_skip.pack(side="left", padx=20)
        self.btn_next = tk.Button(f_foot, text="TIẾP TỤC >>", command=self.on_click_next, bg="#FF9800", fg="white", font=("Arial",12,"bold")); self.btn_next.pack(side="right", padx=20)

    # --- LOGIC FLOW ---
    def load_step_logic(self):
        if self.current_step_index < len(self.selected_queue):
            part = self.selected_queue[self.current_step_index]
            self.lbl_status.config(text=f"ĐANG SOẠN: {part.upper()} ({self.current_step_index+1}/{len(self.selected_queue)})")
            self.entry_kw.delete(0, tk.END); self.listbox.delete(0, tk.END)
            self.search_results = []; self.current_parsed_song = {}
            for c in self.scroll_frame.winfo_children(): c.destroy()
            self.entry_kw.focus(); self.entry_kw.config(state="normal")
            self.btn_search.config(state="normal"); self.btn_fetch.config(state="normal")
            self.btn_next.config(text=f"XÁC NHẬN '{part}' & TIẾP >>", command=self.on_click_next, bg="#FF9800")
            self.btn_skip.pack(side="left", padx=20)
        else:
            self.lbl_status.config(text="HOÀN TẤT!", fg="green")
            self.btn_next.config(text="💾 XUẤT FILE PPTX", command=self.generate_final_pptx, bg="#4CAF50")
            self.btn_skip.pack_forget()
            self.entry_kw.config(state="disabled"); self.btn_search.config(state="disabled"); self.btn_fetch.config(state="disabled")

    def on_click_next(self):
        part = self.selected_queue[self.current_step_index]
        data = []
        for k, item in self.current_parsed_song.items():
            if self.checkbox_vars.get(k) and self.checkbox_vars[k].get():
                data.append({"header": k, "label": item['label'], "content": item['content']})
        if data: self.collected_songs[part] = data
        else:
            if not messagebox.askyesno("Xác nhận", f"Bỏ qua phần '{part}'?"): return
        self.current_step_index += 1
        self.load_step_logic()

    def on_click_skip(self):
        self.current_step_index += 1
        self.load_step_logic()

    # --- SEARCH LOGIC ---
    def on_click_search(self):
        kw = self.entry_kw.get().strip()
        if not kw: return
        if not self.driver: messagebox.showinfo("Lỗi", "Chrome đang mở..."); return
        
        self.root.config(cursor="watch"); self.btn_search.config(state="disabled"); self.entry_kw.config(state="disabled")
        
        if self.is_first_search:
            messagebox.showinfo("Lưu ý Lần Đầu", "Sắp mở Google.\nHãy giải Captcha (nếu có) trên Chrome.\nSau đó chờ phần mềm tự lấy kết quả.\n\nBấm OK để bắt đầu.")
            
        self.is_searching = True
        threading.Thread(target=self._search_polling_worker, args=(kw,), daemon=True).start()

    def _search_polling_worker(self, keyword):
        try:
            self.driver.get(f"https://www.google.com/search?q=site:thanhcavietnam.net {keyword}")
            if self.is_first_search: self.driver.maximize_window(); self.is_first_search = False
            else: self.driver.minimize_window()
            
            max_retries = 120; found_links = []
            for _ in range(max_retries):
                if not self.is_searching: break
                try:
                    elements = self.driver.find_elements(By.CSS_SELECTOR, "a")
                    temp = []
                    for l in elements:
                        h = l.get_attribute("href"); t = l.text
                        if h and "thanhcavietnam.net" in h and "google" not in h and t.strip():
                            temp.append({"title": t, "href": h})
                    if temp: found_links = temp; break 
                except: pass
                time.sleep(0.5)
            
            self.search_results = []
            seen = set()
            for item in found_links:
                if item['href'] not in seen: self.search_results.append(item); seen.add(item['href'])
            self.root.after(0, self._finish_search)
        except: self.root.after(0, self._finish_search)

    def _finish_search(self):
        self.listbox.delete(0, tk.END)
        for item in self.search_results: self.listbox.insert(tk.END, item['title'])
        self.root.config(cursor=""); self.btn_search.config(state="normal"); self.entry_kw.config(state="normal"); self.entry_kw.focus()
        try: self.driver.minimize_window()
        except: pass
        self.root.deiconify(); self.root.lift(); self.root.focus_force()

    # --- FETCH & CLEAN LOGIC (FIX NO CONTENT + HEADER CUT + FOOTER CUT) ---
    def on_click_fetch(self):
        sel = self.listbox.curselection()
        if not sel: return
        url = self.search_results[sel[0]]['href']
        self.root.config(cursor="watch")
        threading.Thread(target=self._fetch_worker, args=(url,), daemon=True).start()

    def _fetch_worker(self, url):
        try:
            headers = {'User-Agent': 'Mozilla/5.0'}
            resp = requests.get(url, headers=headers)
            soup = BeautifulSoup(resp.text, 'html.parser')
            
            # 1. Dọn rác DOM an toàn (Không xóa table/div chứa nội dung)
            for comment in soup.find_all(string=lambda text: isinstance(text, Comment)): comment.extract()
            # Xóa các thẻ script/style nhưng GIỮ LẠI table và div
            for tag in soup.find_all(["script", "style", "iframe", "object", "fieldset"]): tag.decompose()
            
            # Tìm vùng nội dung
            content_div = soup.find("div", id=re.compile(r"^post_message_"))
            if content_div:
                q = content_div.find("blockquote")
                # Lấy text, giữ lại xuống dòng
                txt = q.get_text(separator='\n') if q else content_div.get_text(separator='\n')
                
                # --- XỬ LÝ LỌC VĂN BẢN ---
                lines = txt.split('\n')
                final_lines = []
                
                # Danh sách từ khóa
                header_keywords = ["SÁNG TÁC", "TÁC GIẢ", "LỜI:", "NHẠC:", "THƠ:", "IMPRIMATUR", "Ý:"]
                stop_keywords = ["VIEW MORE", "CÁC CHỦ ĐỀ", "COPYRIGHT", "POWERED BY"]
                # PDF/MP3 ở cuối bài thì cắt, ở đầu bài thì bỏ qua dòng
                skip_keywords = ["PDF", "MP3", "ENCORE", "LINK DOWNLOAD", "BẢN NHẠC", "FILE"]

                # 2. CẮT ĐẦU (Header Cut)
                # Tìm vị trí metadata cuối cùng ở phần đầu
                start_index = 0
                for i, line in enumerate(lines):
                    if any(kw in line.upper() for kw in header_keywords):
                        start_index = i + 1
                lines = lines[start_index:] 

                # 3. DUYỆT VÀ CẮT ĐUÔI
                for line in lines:
                    raw = line.strip()
                    # Quan trọng: Giữ lại dòng trống để logic tách đoạn hoạt động
                    if not raw: 
                        final_lines.append("") 
                        continue
                    
                    upper = raw.upper()
                    
                    # Nếu gặp từ khóa Dừng -> Cắt hết phần sau
                    if any(kw in upper for kw in stop_keywords): break
                    
                    # Nếu gặp link tải (ngắn) -> Bỏ qua dòng này hoặc Dừng luôn nếu nó nằm ở cuối
                    # Để an toàn, với PDF/MP3, ta sẽ bỏ qua dòng này (Skip)
                    is_skip = False
                    for kw in skip_keywords:
                        if kw in upper and len(raw) < 50:
                            is_skip = True; break
                    if is_skip: continue 
                    
                    final_lines.append(raw)
                
                clean_text = "\n".join(final_lines)
                self._parse_text(clean_text)
                self.root.after(0, self._update_parsed_ui)
            else: self.root.after(0, lambda: messagebox.showerror("Lỗi", "Không tìm thấy nội dung"))
        except Exception as e: self.root.after(0, lambda: messagebox.showerror("Lỗi", str(e)))
        finally: self.root.after(0, lambda: self.root.config(cursor=""))

    def _parse_text(self, text):
        self.current_parsed_song = {}
        lines = text.split('\n') # Giữ nguyên dòng trống
        
        p_chorus = re.compile(r'^(ĐK|DK|Đk|Dk|Điệp Khúc|Diep Khuc|Chorus)[:\.\s]?', re.IGNORECASE)
        p_verse = re.compile(r'^(\d+)([\.\)\/\-\:])')
        p_end = re.compile(r'^(CODA|KẾT|KEÁT|FINE)', re.IGNORECASE)

        verse_count = 1
        head = f"Phiên Khúc {verse_count}"
        lbl = f"{verse_count}."
        content = []

        def save(h, l, c):
            if c:
                clean_c = [line for line in c if line.strip()]
                if clean_c:
                    k = h; cnt = 1
                    while k in self.current_parsed_song: k = f"{h} ({cnt})"; cnt+=1
                    self.current_parsed_song[k] = {"label": l, "content": "\n".join(clean_c)}

        for line in lines:
            s_line = line.strip()
            
            # --- LOGIC TÁCH ĐOẠN BẰNG DÒNG TRỐNG (Cho bài không đánh số) ---
            if not s_line:
                if content: # Nếu đang có nội dung -> Lưu đoạn cũ
                    save(head, lbl, content)
                    content = []
                    # Logic chuyển đoạn:
                    # Nếu vừa xong ĐK hoặc Kết -> Đoạn tiếp theo chắc chắn là PK mới
                    # Nếu vừa xong PK -> Đoạn tiếp theo cũng là PK mới (nếu bài viết tách khổ bằng enter)
                    if "Điệp Khúc" in head or "KẾT" in head.upper() or "Phiên Khúc" in head:
                        if "Điệp Khúc" in head or "KẾT" in head.upper():
                            verse_count += 1
                        else:
                            # Nếu đang là PK, gặp dòng trống -> Có thể là PK tiếp theo
                            # Nhưng để tránh nát vụn do xuống dòng lung tung, ta chỉ tăng nếu user
                            # cố tình để dòng trống. Ở đây ta cứ tăng cho bài Ra Khơi.
                            verse_count += 1
                            
                        head = f"Phiên Khúc {verse_count}"
                        lbl = f"{verse_count}."
                continue

            # Kiểm tra Regex Tiêu đề (Ưu tiên cao hơn)
            is_new = False
            m_chorus = p_chorus.match(s_line)
            m_verse = p_verse.match(s_line)
            m_end = p_end.match(s_line)
            start_idx = 0
            
            if m_chorus:
                is_new = True; head = "Điệp Khúc"; lbl = "ĐK:"; start_idx = m_chorus.end()
            elif m_verse:
                is_new = True; num = m_verse.group(1); verse_count = int(num) # Sync số đếm
                head = f"Phiên Khúc {num}"; lbl = f"{num}."; start_idx = m_verse.end()
            elif m_end:
                is_new = True; head = s_line.strip(); lbl = s_line.strip(); start_idx = len(s_line)
            
            if is_new:
                save(head, lbl, content) # Lưu đoạn trước
                content = []
                rem = s_line[start_idx:].strip()
                while rem and rem[0] in [".", ":", " ", ")", "/"]: rem = rem[1:].strip()
                if rem: content.append(rem)
            else:
                content.append(s_line)
        
        save(head, lbl, content) # Lưu đoạn cuối

    def _update_parsed_ui(self):
        for c in self.scroll_frame.winfo_children(): c.destroy()
        self.checkbox_vars = {}
        for k, v in self.current_parsed_song.items():
            var = tk.BooleanVar(value=True); self.checkbox_vars[k] = var
            r = tk.Frame(self.scroll_frame); r.pack(fill="x", pady=2)
            tk.Checkbutton(r, text=k, variable=var, font=("Arial",10,"bold"), fg="#0D47A1").pack(side="left")
            tk.Label(r, text=f"[{v['label']}] {v['content'][:40]}...", fg="gray").pack(side="left")

    # --- GENERATE PPTX (FIX TRÙNG SLIDE) ---
    def generate_final_pptx(self):
        path = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint", "*.pptx")])
        if not path: return
        prs = Presentation()
        prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

        def add_black():
            s = prs.slides.add_slide(prs.slide_layouts[6])
            s.background.fill.solid(); s.background.fill.fore_color.rgb = RGBColor(0,0,0)

        def has_song(key): return key in self.collected_songs
        def has_read(key): return bool(self.ktcgkpv_data.get(key))

        def add_song(key):
            if not has_song(key): return
            add_black()
            sections = self.collected_songs[key]
            
            chorus = next((s for s in sections if "Điệp Khúc" in s['header'] or "ĐK" in s['header']), None)
            has_verses = any("Phiên Khúc" in s['header'] for s in sections)

            for s in sections:
                # --- LOGIC CHỐNG TRÙNG ---
                # Nếu bài hát có Phiên Khúc, thì ĐK đã được kẹp sau mỗi PK rồi.
                # Nên ta BỎ QUA việc tạo slide cho ĐK khi nó xuất hiện lẻ trong vòng lặp này.
                if has_verses and ("Điệp Khúc" in s['header'] or "ĐK" in s['header']):
                    continue
                
                self.create_slide(prs, key.upper(), s['label'], s['content'])
                
                # Tự động kẹp ĐK sau PK
                if "Phiên Khúc" in s['header'] and chorus:
                    self.create_slide(prs, key.upper(), chorus['label'], chorus['content'])

        def add_read(name):
            txt = self.ktcgkpv_data.get(name, "")
            if not txt: return
            add_black()
            self.create_slide(prs, name.upper(), "", txt)

        if has_song("Đức Mẹ"): add_song("Đức Mẹ")
        if has_song("Nhập Lễ"): add_song("Nhập Lễ")
        if has_read("Ca Nhập Lễ"): add_read("Ca Nhập Lễ")

        if has_read("Đáp Ca") or has_read("Tung Hô Tin Mừng"):
            add_black()
            if has_read("Đáp Ca"): self.create_slide(prs, "ĐÁP CA", "", self.ktcgkpv_data["Đáp Ca"])
            if has_read("Tung Hô Tin Mừng"): self.create_slide(prs, "TUNG HÔ TIN MỪNG", "", self.ktcgkpv_data["Tung Hô Tin Mừng"])

        if has_song("Dâng Lễ"): add_song("Dâng Lễ")
        if has_read("Ca Hiệp Lễ"): add_read("Ca Hiệp Lễ")
        if has_song("Hiệp Lễ 1"): add_song("Hiệp Lễ 1")
        if has_song("Hiệp Lễ 2"): add_song("Hiệp Lễ 2")
        if has_song("Kết Lễ"): add_song("Kết Lễ")

        try:
            prs.save(path)
            messagebox.showinfo("Xong", f"Đã lưu: {path}")
            os.startfile(path)
        except Exception as e: messagebox.showerror("Lỗi", str(e))

    def create_slide(self, prs, title, label, content):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid(); slide.background.fill.fore_color.rgb = RGBColor(0, 32, 96)
        sw, sh = prs.slide_width, prs.slide_height

        tb = slide.shapes.add_textbox(Inches(0), Inches(0), sw, Inches(1.4))
        p = tb.text_frame.paragraphs[0]
        p.text = title; p.alignment = PP_ALIGN.CENTER
        p.font.name = "Times New Roman"; p.font.size = Pt(72); p.font.bold = True; p.font.color.rgb = RGBColor(255, 255, 0)

        ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0), Inches(1.5), sw, Inches(1.5))
        ln.line.color.rgb = RGBColor(255, 255, 0); ln.line.width = Pt(3)

        tb2 = slide.shapes.add_textbox(Inches(0), Inches(1.6), sw, sh - Inches(2.0))
        tf = tb2.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
        
        raw = " ".join([l.strip() for l in content.split('\n') if l.strip()])
        length = len(raw)
        lines_54 = math.ceil(length / 32)
        lines_44 = math.ceil(length / 42) 
        if lines_54 <= 6: size = 54
        elif lines_44 <= 9: size = 44
        elif length < 700: size = 36
        else: size = 28

        p2 = tf.paragraphs[0]; p2.alignment = PP_ALIGN.JUSTIFY; p2.line_spacing = 1.1
        if label:
            r = p2.add_run(); r.text = label + " "
            r.font.name = "Times New Roman"; r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = RGBColor(255, 0, 0)
        r2 = p2.add_run(); r2.text = raw
        r2.font.name = "Times New Roman"; r2.font.size = Pt(size); r2.font.bold = True; r2.font.color.rgb = RGBColor(255, 255, 255)

if __name__ == "__main__":
    root = tk.Tk()
    app = MassSlideApp(root)
    root.mainloop()